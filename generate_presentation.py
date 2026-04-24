#!/usr/bin/env python3
"""
Generate a 15-slide presentation:
"Teaching Math in the Browser: Pedagogical Philosophy of math-tools.org"
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Color palette
DARK_NAVY = RGBColor(26, 39, 68)      # #1a2744
GOLD = RGBColor(212, 175, 55)         # #D4AF37
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(230, 230, 235)
MED_GRAY = RGBColor(100, 100, 115)

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    slide1 = add_title_slide(prs,
        "Teaching Math in the Browser",
        "Pedagogical Philosophy of math-tools.org",
        "Shelvean Kapita, Ph.D.\nTexas A&M University")

    # Slide 2: The Problem
    slide2 = add_content_slide(prs, "The Problem")
    add_bullet_points(slide2,
        ["Textbook pictures are static",
         "Require imagination from readers",
         "One example ≠ deep understanding",
         "No feedback on \"what if I change...?\""],
        DARK_NAVY, MED_GRAY, WHITE)

    # Slide 3: The Solution
    slide3 = add_content_slide(prs, "The Solution: Interactive Visualization")
    add_bullet_points(slide3,
        ["Real-time parameter adjustment",
         "Generate lots of examples instantly",
         "See abstract concepts in motion",
         "Feedback loop: change → visualize → understand"],
        DARK_NAVY, MED_GRAY, WHITE)

    # Slide 4: Scope & Reach
    slide4 = add_content_slide(prs, "Scope: 70+ Interactive Tools")
    left_box = slide4.shapes.add_shape(1, Inches(0.5), Inches(1.8), Inches(4.5), Inches(5))
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = LIGHT_GRAY
    left_box.line.color.rgb = GOLD
    left_box.line.width = Pt(2)
    tf = left_box.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.2)
    tf.margin_left = Inches(0.2)
    p = tf.paragraphs[0]
    p.text = "116,805 page views\n7,490 active users\n115 countries\nMIT licensed\nWCAG 2.1 AA"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DARK_NAVY

    right_box = slide4.shapes.add_shape(1, Inches(5.2), Inches(1.8), Inches(4.3), Inches(5))
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = LIGHT_GRAY
    right_box.line.color.rgb = GOLD
    right_box.line.width = Pt(2)
    tf2 = right_box.text_frame
    tf2.word_wrap = True
    tf2.margin_top = Inches(0.2)
    tf2.margin_left = Inches(0.2)
    p2 = tf2.paragraphs[0]
    p2.text = "Since Jan 2025\nGlobal audience\nFree, no installs\nBrowser-native"
    p2.font.size = Pt(18)
    p2.font.bold = True
    p2.font.color.rgb = DARK_NAVY

    # Slide 5: Why it Matters
    slide5 = add_content_slide(prs, "Why It Matters: Math 308")
    add_bullet_points(slide5,
        ["Engineering students need tangible connections",
         "Abstract PDEs → real mass-spring, heat flow, pendula",
         "Nonlinear dynamics → control, stability, design",
         "Interactive tools = active learning"],
        DARK_NAVY, MED_GRAY, WHITE)

    # Slide 6: Physical Examples 1
    slide6 = add_content_slide(prs, "Physical Examples: Oscillations")
    add_bullet_points(slide6,
        ["Mass-spring (horizontal & vertical)",
         "Coupled oscillators → normal modes",
         "Forcing & resonance → real frequency response",
         "Damping (underdamped → overdamped)"],
        DARK_NAVY, MED_GRAY, WHITE)
    add_tools_box(slide6, "massspring.html · coupled.html · forcing.html", GOLD)

    # Slide 7: Physical Examples 2
    slide7 = add_content_slide(prs, "Physical Examples: Pendula")
    add_bullet_points(slide7,
        ["Simple pendulum → direction fields",
         "Elastic pendulum → stretching meets gravity",
         "Multi-pendulum → chaos emerges",
         "Pendulum wave → beats & synchronization"],
        DARK_NAVY, MED_GRAY, WHITE)
    add_tools_box(slide7, "pend.html · elastic_pendulum.html · pendulum_wave.html", GOLD)

    # Slide 8: Dynamical Systems
    slide8 = add_content_slide(prs, "Dynamical Systems & Bifurcation")
    add_bullet_points(slide8,
        ["Phase portraits: nodes, spirals, saddles",
         "Bifurcation: period-doubling cascade",
         "Strange attractors: Lorenz, Rössler, Van der Pol",
         "Nonlinear dynamics as unifying theme"],
        DARK_NAVY, MED_GRAY, WHITE)
    add_tools_box(slide8, "bifurcation.html · lorenzsystem.html · vanderpol.html · duffing.html", GOLD)

    # Slide 9: Linear Algebra
    slide9 = add_content_slide(prs, "Linear Algebra: Geometry Over Computation")
    add_bullet_points(slide9,
        ["RREF: step-by-step with all pivot operations",
         "SVD: how transformations stretch & rotate spaces",
         "Eigenvectors: invariant directions, diagonalization",
         "Determinants: volume scaling, inverses"],
        DARK_NAVY, MED_GRAY, WHITE)
    add_tools_box(slide9, "rref.html · svd2d.html · eigenvalues-eigenvectors.html · detviz.html", GOLD)

    # Slide 10: Why JavaScript?
    slide10 = add_content_slide(prs, "Why JavaScript? The Web's Native Language")
    add_bullet_points(slide10,
        ["Every browser comes with a programming language",
         "No plugins, no installation, no dependencies",
         "Runs on phones, tablets, laptops, PCs",
         "Instant updates (refresh = latest version)"],
        DARK_NAVY, MED_GRAY, WHITE)

    # Slide 11: Tech Stack
    slide11 = add_content_slide(prs, "Tech Stack: CDN-Powered Libraries")
    add_bullet_points(slide11,
        ["MathJax 3: math rendering + screen reader support",
         "KaTeX: fast LaTeX rendering",
         "D3.js, Plotly.js: interactive plotting",
         "math.js: symbolic computation",
         "All loaded from CDNs—no build step, no local install"],
        DARK_NAVY, MED_GRAY, WHITE)

    # Slide 12: Workflow Evolution
    slide12 = add_content_slide(prs, "Workflow: 2025 → 2026 with Claude Code")
    left_col = slide12.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4), Inches(5))
    tf_l = left_col.text_frame
    tf_l.word_wrap = True
    p_l = tf_l.paragraphs[0]
    p_l.text = "2025\n"
    p_l.font.size = Pt(20)
    p_l.font.bold = True
    p_l.font.color.rgb = GOLD
    p_l = tf_l.add_paragraph()
    p_l.text = "Prompt ChatGPT/Grok\n→ Copy JS\n→ Notepad++ test\n→ Iterate\n→ Git commit"
    p_l.font.size = Pt(14)
    p_l.font.color.rgb = WHITE
    p_l.space_before = Pt(6)

    right_col = slide12.shapes.add_textbox(Inches(5.2), Inches(1.8), Inches(4), Inches(5))
    tf_r = right_col.text_frame
    tf_r.word_wrap = True
    p_r = tf_r.paragraphs[0]
    p_r.text = "2026\n"
    p_r.font.size = Pt(20)
    p_r.font.bold = True
    p_r.font.color.rgb = GOLD
    p_r = tf_r.add_paragraph()
    p_r.text = "Prompt Claude Code\n→ Agents commit to GitHub\n→ Change ALL files at once\n→ Accessibility sweep, design updates\n→ Iterate & re-commit"
    p_r.font.size = Pt(14)
    p_r.font.color.rgb = WHITE
    p_r.space_before = Pt(6)

    # Slide 13: Accessibility
    slide13 = add_content_slide(prs, "Accessibility: WCAG 2.1 AA Across All 70+ Tools")
    add_bullet_points(slide13,
        ["ARIA labels on every interactive element",
         "MathJax with assistive MathML (screen readers read math)",
         "Keyboard navigation: no mouse required",
         "High-contrast mode toggle, color-blind safe"],
        DARK_NAVY, MED_GRAY, WHITE)

    # Slide 14: Global Reach & Testimonials
    slide14 = add_content_slide(prs, "Global Reach & Impact")
    add_bullet_points(slide14,
        ["\"I think your work is just wonderful!\"—Gilbert Strang, MIT",
         "\"Beautiful! Very elegant.\"—Steven Strogatz, Cornell",
         "Provost OER Grant for Math 308 question bank (2026)",
         "Freely available: open education for global equity"],
        DARK_NAVY, MED_GRAY, WHITE)

    # Slide 15: Closing
    slide15 = add_content_slide(prs, "Bring Scientific Computing to the Web")
    center_box = slide15.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(3.5))
    tf_c = center_box.text_frame
    tf_c.word_wrap = True
    p_c = tf_c.paragraphs[0]
    p_c.text = "math-tools.org"
    p_c.font.size = Pt(44)
    p_c.font.bold = True
    p_c.font.color.rgb = GOLD
    p_c.alignment = PP_ALIGN.CENTER
    p_c = tf_c.add_paragraph()
    p_c.text = "\n70+ interactive visualizations\nfor Differential Equations, Linear Algebra & more"
    p_c.font.size = Pt(16)
    p_c.font.color.rgb = WHITE
    p_c.alignment = PP_ALIGN.CENTER
    p_c.space_before = Pt(12)

    return prs

def add_title_slide(prs, title, subtitle, author):
    """Create a title slide with consistent styling."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_NAVY

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(9), Inches(1))
    tf = subtitle_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Author
    author_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1.5))
    tf = author_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = author
    p.font.size = Pt(16)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title):
    """Create a content slide with title and dark background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_NAVY

    # Gold line under title
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.2), Inches(9), Inches(0))
    line.line.color.rgb = GOLD
    line.line.width = Pt(3)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = GOLD

    return slide

def add_bullet_points(slide, bullets, bg_color, bullet_color, text_color):
    """Add bullet points to a content slide."""
    text_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
    tf = text_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(20)
        p.font.color.rgb = text_color
        p.space_before = Pt(10)
        p.space_after = Pt(10)

def add_tools_box(slide, tools_text, box_color):
    """Add a small box with example tool names at the bottom."""
    box = slide.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(8.4), Inches(0.8))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Examples: " + tools_text
    p.font.size = Pt(11)
    p.font.italic = True
    p.font.color.rgb = box_color

if __name__ == "__main__":
    prs = create_presentation()
    prs.save("/home/user/math-tools/Math-Tools-Pedagogical-Philosophy.pptx")
    print("✓ Generated: Math-Tools-Pedagogical-Philosophy.pptx")
    print("  15 slides, ~10 minutes")
    print("  Dark navy + gold aesthetic")
    print("  Ready for presentation!")
